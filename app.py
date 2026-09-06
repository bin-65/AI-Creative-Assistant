import streamlit as st
from google import genai
import pypdf
import docx
import pptx
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
import io
import time

# Page configuration
st.set_page_config(
    page_title="Usman AI Studio", 
    page_icon="⚡", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom Styling (PowerPoint Design Gallery Card Styling)
st.markdown("""
    <style>
    .stApp { background-color: #f8fafc; color: #1e293b; }
    .header-box {
        background: linear-gradient(135deg, #e0f2fe 0%, #bae6fd 100%);
        padding: 24px; border-radius: 12px; border: 1px solid #7dd3fc;
        margin-bottom: 20px; box-shadow: 0 4px 12px rgba(14, 165, 233, 0.12);
    }
    .main-title { color: #0369a1; font-size: 28px; font-weight: 800; margin: 0; }
    .sub-title { color: #0284c7; font-size: 14px; margin-top: 4px; font-weight: 500; }
    
    .stButton>button {
        background: linear-gradient(135deg, #0284c7 0%, #0369a1 100%) !important;
        color: #ffffff !important; border: none !important;
        border-radius: 8px !important; font-weight: 600 !important;
        width: 100% !important; padding: 6px 12px !important; margin-top: 4px !important;
    }
    
    /* PPT Template Preview Cards */
    .ppt-card {
        border: 2px solid #cbd5e1; border-radius: 8px; padding: 12px;
        height: 120px; display: flex; flex-direction: column; justify-content: space-between;
        margin-bottom: 6px; box-shadow: 0 2px 5px rgba(0,0,0,0.05);
    }
    .palette-dots { display: flex; gap: 4px; margin-top: 6px; }
    .dot { width: 12px; height: 12px; border-radius: 50%; display: inline-block; }
    </style>
""", unsafe_allow_html=True)

# Main Top Header Bar
st.markdown("""
    <div class="header-box">
        <div class="main-title">⚡ Usman AI Studio</div>
        <div class="sub-title">Powered by Google Gemini AI</div>
    </div>
""", unsafe_allow_html=True)

# Helper Export Functions
def generate_pdf_bytes(title, text_content):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 14)
    pdf.cell(0, 10, title.encode('latin-1', 'replace').decode('latin-1'), ln=True)
    pdf.set_font("Arial", size=10)
    pdf.ln(5)
    clean_text = text_content.encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 7, clean_text)
    return bytes(pdf.output())

def generate_docx_bytes(title, text_content):
    doc = Document()
    doc.add_heading(title, 0)
    for paragraph in text_content.split("\n\n"):
        if paragraph.strip():
            doc.add_paragraph(paragraph.strip())
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# Advanced PPT Engine generating MS PowerPoint Native Themes
def generate_pptx_bytes(title, text_content, theme_name="Facet Modern"):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # 10 PowerPoint Themes Styling Dictionary
    PPT_THEMES = {
        "Facet Modern": {"bg": RGBColor(255, 255, 255), "primary": RGBColor(3, 169, 244), "title_color": RGBColor(1, 87, 155), "text": RGBColor(33, 33, 33), "font": "Segoe UI", "layout": "top_bar"},
        "Integral Dark": {"bg": RGBColor(33, 33, 33), "primary": RGBColor(255, 112, 67), "title_color": RGBColor(255, 171, 145), "text": RGBColor(245, 245, 245), "font": "Trebuchet MS", "layout": "full_dark"},
        "Ion Purple": {"bg": RGBColor(250, 245, 255), "primary": RGBColor(147, 51, 234), "title_color": RGBColor(88, 28, 135), "text": RGBColor(58, 12, 90), "font": "Calibri", "layout": "side_stripe"},
        "Organic Emerald": {"bg": RGBColor(240, 253, 244), "primary": RGBColor(16, 185, 129), "title_color": RGBColor(6, 78, 59), "text": RGBColor(20, 83, 45), "font": "Georgia", "layout": "bottom_accent"},
        "Slice Crimson": {"bg": RGBColor(255, 255, 255), "primary": RGBColor(225, 29, 72), "title_color": RGBColor(136, 19, 55), "text": RGBColor(30, 41, 59), "font": "Arial Black", "layout": "top_bar"},
        "Retrospect Gold": {"bg": RGBColor(254, 252, 232), "primary": RGBColor(202, 138, 4), "title_color": RGBColor(113, 63, 18), "text": RGBColor(66, 32, 6), "font": "Garamond", "layout": "side_stripe"},
        "Slate Minimal": {"bg": RGBColor(248, 250, 252), "primary": RGBColor(71, 85, 105), "title_color": RGBColor(15, 23, 42), "text": RGBColor(51, 65, 85), "font": "Calibri", "layout": "bottom_accent"},
        "Vapor Neon": {"bg": RGBColor(15, 23, 42), "primary": RGBColor(6, 182, 212), "title_color": RGBColor(165, 243, 252), "text": RGBColor(226, 232, 240), "font": "Consolas", "layout": "full_dark"},
        "Warm Ochre": {"bg": RGBColor(255, 247, 237), "primary": RGBColor(234, 88, 12), "title_color": RGBColor(154, 52, 18), "text": RGBColor(67, 20, 7), "font": "Tahoma", "layout": "top_bar"},
        "Teal Tech": {"bg": RGBColor(240, 253, 250), "primary": RGBColor(13, 148, 136), "title_color": RGBColor(19, 78, 74), "text": RGBColor(15, 118, 110), "font": "Segoe UI", "layout": "side_stripe"},
    }

    thm = PPT_THEMES.get(theme_name, PPT_THEMES["Facet Modern"])

    # Slide 1: Main Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    s1_bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    s1_bg.fill.solid()
    s1_bg.fill.fore_color.rgb = thm["bg"]
    s1_bg.line.fill.background()

    # Title Card Accent
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = thm["primary"]
    card.line.fill.background()

    tf1 = card.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = title
    p1.font.name = thm["font"]
    p1.font.bold = True
    p1.font.size = Pt(32)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = f"\nPowerPoint Theme: {theme_name}\nGenerated by Usman AI Studio"
    p2.font.name = thm["font"]
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(245, 245, 245)
    p2.alignment = PP_ALIGN.CENTER

    # Content Slides
    paragraphs = [p.strip() for p in text_content.split("\n\n") if p.strip()]
    chunk_size = 3

    for i in range(0, len(paragraphs), chunk_size):
        chunk = paragraphs[i:i + chunk_size]
        slide = prs.slides.add_slide(blank_layout)

        # Slide Background
        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = thm["bg"]
        s_bg.line.fill.background()

        # Layout Geometry
        if thm["layout"] == "top_bar":
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = thm["primary"]
            bar.line.fill.background()
            
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{title} - Slide {(i // chunk_size) + 1}"
            p.font.name = thm["font"]
            p.font.bold = True
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)

        elif thm["layout"] == "side_stripe":
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = thm["primary"]
            stripe.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.8), Inches(0.8))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{title} - Slide {(i // chunk_size) + 1}"
            p.font.name = thm["font"]
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.color.rgb = thm["title_color"]

        else: # full_dark / bottom_accent
            b_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(0.08))
            b_line.fill.solid()
            b_line.fill.fore_color.rgb = thm["primary"]
            b_line.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{title} - Slide {(i // chunk_size) + 1}"
            p.font.name = thm["font"]
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.color.rgb = thm["title_color"]

        # Content Text Box
        cnt_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
        tf = cnt_box.text_frame
        tf.word_wrap = True

        for idx, para in enumerate(chunk):
            p_cnt = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p_cnt.text = f"•  {para}"
            p_cnt.font.name = thm["font"]
            p_cnt.font.size = Pt(15)
            p_cnt.font.color.rgb = thm["text"]
            p_cnt.space_after = Pt(14)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

# Gemini AI Engine
def call_gemini_ai(prompt_text):
    api_key = st.secrets.get("GEMINI_API_KEY", "")
    if not api_key:
        raise Exception("GEMINI_API_KEY Missing in Streamlit Secrets!")
    
    clean_key = str(api_key).strip().strip('"').strip("'")
    client = genai.Client(api_key=clean_key)
    models_to_try = ['gemini-3.6-flash', 'gemini-2.5-flash']
    
    for model_name in models_to_try:
        for attempt in range(3):
            try:
                response = client.models.generate_content(
                    model=model_name,
                    contents=prompt_text,
                )
                return response.text
            except Exception as e:
                error_msg = str(e)
                if "503" in error_msg or "UNAVAILABLE" in error_msg:
                    time.sleep(1.5)
                    continue
                else:
                    break
                    
    raise Exception("Server is busy. Please try clicking generate again in a few seconds.")

# Sidebar
st.sidebar.markdown("### ⚙️ System Status")
st.sidebar.success("✨ **Usman AI Studio Active**")
st.sidebar.markdown("---")

# Main Tabs Navigation
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📊 Presentation Studio",
    "✍️ Content Creator", 
    "🌐 Translator", 
    "⚡ Smart AI Workspace",
    "📚 Academic Writer",
    "📑 Advanced Doc Hub"
])

def render_output_section(result_text, doc_title, key_prefix, show_ppt=True, theme_name="Facet Modern"):
    st.markdown("---")
    st.success("🎉 Generated Successfully!")
    
    try:
        pdf_data = generate_pdf_bytes(doc_title, result_text)
        docx_data = generate_docx_bytes(doc_title, result_text)
        pptx_data = generate_pptx_bytes(doc_title, result_text, theme_name)
        
        if show_ppt:
            c1, c2, c3 = st.columns([1, 1, 1])
            with c1: st.download_button("📥 Download PDF (.pdf)", pdf_data, f"{doc_title}.pdf", "application/pdf", use_container_width=True, key=f"{key_prefix}_pdf")
            with c2: st.download_button("📄 Download Word (.docx)", docx_data, f"{doc_title}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True, key=f"{key_prefix}_docx")
            with c3: st.download_button("📊 Download Presentation (.pptx)", pptx_data, f"{doc_title}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True, key=f"{key_prefix}_pptx")
        else:
            c1, c2 = st.columns([1, 1])
            with c1: st.download_button("📥 Download PDF (.pdf)", pdf_data, f"{doc_title}.pdf", "application/pdf", use_container_width=True, key=f"{key_prefix}_pdf")
            with c2: st.download_button("📄 Download Word (.docx)", docx_data, f"{doc_title}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True, key=f"{key_prefix}_docx")
    except Exception as export_err:
        st.error(f"Error generating download files: {export_err}")
        
    st.markdown("### 📄 Result Output:")
    st.markdown(result_text)

# TAB 1: PRESENTATION STUDIO (MS POWERPOINT DESIGN GALLERY)
with tab1:
    st.markdown("## Create Slides with AI & Edit with Full Control")
    ppt_mode = st.radio("Creation Mode", ["✨ From Topic", "📄 From Document", "📑 From Outline"], horizontal=True)
    st.markdown("---")
    
    if "From Topic" in ppt_mode:
        topic_input = st.text_area("Slide Topic / Title", value="Four stroke engine", height=100)
        
        c1, c2, c3 = st.columns([2, 2, 2])
        with c1: use_search = st.checkbox("🔍 Enable Web Search", value=True)
        with c2: deep_think = st.checkbox("🧠 DeepThink Reasoning", value=True)
        with c3: slide_count = st.slider("Target Slide Count", 5, 20, 8)

        st.markdown("##### Quick Sample Topics:")
        q1, q2, q3, q4 = st.columns(4)
        if q1.button("Health Management"): topic_input = "Health Management Program"
        if q2.button("Summary Report"): topic_input = "Summary Report"
        if q3.button("Mental Health"): topic_input = "Mental Health Report"
        if q4.button("Data Analysis"): topic_input = "Data Analysis Report"

        st.markdown("---")
        st.markdown("### 🎨 Select PowerPoint Slide Theme (10 Design Templates)")
        
        # 10 MS PowerPoint Gallery Themes
        ppt_themes_gallery = [
            ("Facet Modern", "#ffffff", "#03a9f4", ["#01579b", "#03a9f4", "#81d4fa", "#b3e5fc"]),
            ("Integral Dark", "#212121", "#ff7043", ["#212121", "#ff7043", "#ffab91", "#d84315"]),
            ("Ion Purple", "#faf5ff", "#9333ea", ["#faf5ff", "#9333ea", "#c084fc", "#581c87"]),
            ("Organic Emerald", "#f0fdf4", "#10b981", ["#f0fdf4", "#10b981", "#6ee7b7", "#064e3b"]),
            ("Slice Crimson", "#ffffff", "#e11d48", ["#ffffff", "#e11d48", "#fda4af", "#881337"]),
            ("Retrospect Gold", "#fefce8", "#ca8a04", ["#fefce8", "#ca8a04", "#fde047", "#713f12"]),
            ("Slate Minimal", "#f8fafc", "#475569", ["#f8fafc", "#475569", "#94a3b8", "#0f172a"]),
            ("Vapor Neon", "#0f172a", "#06b6d4", ["#0f172a", "#06b6d4", "#67e8f9", "#164e63"]),
            ("Warm Ochre", "#fff7ed", "#ea580c", ["#fff7ed", "#ea580c", "#fdba74", "#7c2d12"]),
            ("Teal Tech", "#f0fdfa", "#0d9488", ["#f0fdfa", "#0d9488", "#5eead4", "#134e4a"]),
        ]

        # Gallery Cards Rendering (5 Columns x 2 Rows)
        row1 = st.columns(5)
        for idx, (t_name, bg_c, pri_c, palette) in enumerate(ppt_themes_gallery[:5]):
            with row1[idx]:
                dots_html = "".join([f'<span class="dot" style="background-color:{c};"></span>' for c in palette])
                st.markdown(f"""
                    <div class="ppt-card" style="background-color:{bg_c}; border-top: 4px solid {pri_c};">
                        <div style="font-size:22px; font-weight:800; color:{pri_c};">Aa</div>
                        <div style="font-size:12px; font-weight:600; color:#334155;">{t_name}</div>
                        <div class="palette-dots">{dots_html}</div>
                    </div>
                """, unsafe_allow_html=True)
                if st.button(f"Use {t_name}", key=f"thm_{idx}"):
                    st.session_state["active_theme"] = t_name

        row2 = st.columns(5)
        for idx, (t_name, bg_c, pri_c, palette) in enumerate(ppt_themes_gallery[5:]):
            with row2[idx]:
                dots_html = "".join([f'<span class="dot" style="background-color:{c};"></span>' for c in palette])
                st.markdown(f"""
                    <div class="ppt-card" style="background-color:{bg_c}; border-top: 4px solid {pri_c};">
                        <div style="font-size:22px; font-weight:800; color:{pri_c};">Aa</div>
                        <div style="font-size:12px; font-weight:600; color:#334155;">{t_name}</div>
                        <div class="palette-dots">{dots_html}</div>
                    </div>
                """, unsafe_allow_html=True)
                if st.button(f"Use {t_name}", key=f"thm_{idx+5}"):
                    st.session_state["active_theme"] = t_name

        current_thm = st.session_state.get("active_theme", "Facet Modern")
        st.info(f"Selected PowerPoint Theme: **{current_thm}**")

        if st.button("🚀 Generate Presentation Slides", key="ppt_gen_btn"):
            if not topic_input.strip():
                st.warning("⚠️ Topic field empty.")
            else:
                try:
                    with st.spinner(f"Creating Deck with Theme '{current_thm}'..."):
                        ppt_prompt = f"Create a structured {slide_count}-slide presentation deck on '{topic_input}'. Design Theme: {current_thm}. Format slide by slide with titles and detailed bullet points."
                        res_text = call_gemini_ai(ppt_prompt)
                    render_output_section(res_text, f"{topic_input}_Slides", "tab_ppt", show_ppt=True, theme_name=current_thm)
                except Exception as e:
                    st.error(f"Execution Error: {e}")

    elif "From Document" in ppt_mode:
        ppt_file = st.file_uploader("Upload PDF or Word Document", type=["pdf", "docx", "txt"])
        if st.button("🚀 Convert Document to Presentation", key="ppt_doc_btn") and ppt_file:
            try:
                with st.spinner("Converting Document..."):
                    extracted_text = ""
                    if ppt_file.name.endswith(".pdf"):
                        reader = pypdf.PdfReader(ppt_file)
                        extracted_text = "\n".join([page.extract_text() or "" for page in reader.pages[:10]])
                    elif ppt_file.name.endswith(".docx"):
                        doc = docx.Document(ppt_file)
                        extracted_text = "\n".join([p.text for p in doc.paragraphs])
                    res_text = call_gemini_ai(f"Convert this document to slides:\n\n{extracted_text[:8000]}")
                render_output_section(res_text, "Document_Slides", "tab_ppt_doc", show_ppt=True)
            except Exception as e:
                st.error(f"Execution Error: {e}")

    elif "From Outline" in ppt_mode:
        outline_text = st.text_area("Paste Outline", height=150)
        if st.button("🚀 Build Slides from Outline", key="ppt_out_btn") and outline_text.strip():
            try:
                with st.spinner("Building Slides..."):
                    res_text = call_gemini_ai(f"Expand outline to slide deck:\n{outline_text}")
                render_output_section(res_text, "Outline_Slides", "tab_ppt_out", show_ppt=True)
            except Exception as e:
                st.error(f"Execution Error: {e}")

# TAB 2: CONTENT ASSISTANT
with tab2:
    st.subheader("✍️ Social Media Content Generator")
    c1, c2 = st.columns(2)
    with c1:
        platform = st.selectbox("Target Platform", ["LinkedIn", "Instagram", "Twitter / X", "Facebook"])
        content_type = st.selectbox("Content Style", ["Informational Post", "Promotional / Ad", "Storytelling"])
    with c2:
        tone = st.selectbox("Tone & Persona", ["Professional", "Casual & Friendly", "Persuasive"])
        target_audience = st.text_input("Target Audience", placeholder="e.g., Engineers, Students")
    
    topic = st.text_area("Core Brief / Topic")
    if st.button("🚀 Generate Content", key="tab2_btn") and topic and target_audience:
        try:
            res_text = call_gemini_ai(f"Platform: {platform}\nType: {content_type}\nTone: {tone}\nAudience: {target_audience}\nTopic: {topic}")
            render_output_section(res_text, f"{platform}_Content", "tab2", show_ppt=False)
        except Exception as e:
            st.error(f"Execution Error: {e}")

# TAB 3: TRANSLATOR
with tab3:
    st.subheader("🌐 Global Multi-Language Translator")
    target_language = st.selectbox("Target Language", ["English", "Urdu", "Arabic", "Hindi", "Spanish", "French", "German"])
    input_text = st.text_area("Source Text", height=150)
    if st.button("🌐 Translate Content", key="tab3_btn") and input_text.strip():
        try:
            res_text = call_gemini_ai(f"Translate to {target_language}:\n\n{input_text}")
            render_output_section(res_text, f"Translation_{target_language}", "tab3", show_ppt=False)
        except Exception as e:
            st.error(f"Execution Error: {e}")

# TAB 4: SMART WORKSPACE
with tab4:
    st.subheader("⚡ Smart File Workspace")
    user_prompt = st.text_area("Analysis Prompt", height=100)
    uploaded_files = st.file_uploader("Upload Documents", type=["pdf", "docx", "pptx", "txt"], accept_multiple_files=True)
    if st.button("⚡ Process Query", key="tab4_btn"):
        try:
            extracted_text = ""
            if uploaded_files:
                for file in uploaded_files:
                    if file.name.endswith(".pdf"):
                        reader = pypdf.PdfReader(file)
                        extracted_text += "\n".join([p.extract_text() or "" for p in reader.pages[:10]])
                    elif file.name.endswith(".docx"):
                        doc = docx.Document(file)
                        extracted_text += "\n".join([p.text for p in doc.paragraphs])
            res_text = call_gemini_ai(f"Extracted Text:\n{extracted_text[:8000]}\nUser Prompt: {user_prompt}")
            render_output_section(res_text, "Workspace_Output", "tab4", show_ppt=True)
        except Exception as e:
            st.error(f"Execution Error: {e}")

# TAB 5: ACADEMIC WRITER
with tab5:
    st.subheader("📚 Academic & Assignment Writer")
    subject_topic = st.text_input("Topic Header")
    academic_level = st.selectbox("Academic Level", ["Undergraduate", "Postgraduate / PhD", "College"])
    if st.button("✨ Generate Assignment", key="tab5_btn") and subject_topic.strip():
        try:
            res_text = call_gemini_ai(f"Write paper on '{subject_topic}' for {academic_level} level.")
            render_output_section(res_text, f"{subject_topic}_Assignment", "tab5", show_ppt=True)
        except Exception as e:
            st.error(f"Execution Error: {e}")

# TAB 6: DOCUMENT HUB
with tab6:
    st.subheader("📑 Advanced Document Hub")
    doc_topic = st.text_input("Document Topic")
    if st.button("✨ Build Document", key="tab6_btn") and doc_topic.strip():
        try:
            res_text = call_gemini_ai(f"Create a detailed document on '{doc_topic}'.")
            render_output_section(res_text, f"{doc_topic}_Document", "tab6", show_ppt=True)
        except Exception as e:
            st.error(f"Execution Error: {e}")
